"""PowerPoint assembly module."""

import logging
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches
from pptx.util import Pt

from ..models import SlideContent
from ..models import SlideDiagram
from ..models import SlideImage

logger = logging.getLogger(__name__)


async def assemble_presentation(
    slides: list[SlideContent],
    diagrams: list[SlideDiagram],
    images: list[SlideImage],
    template_path: Path | None,
    output_path: Path,
) -> Path:
    """Assemble PowerPoint presentation from components.

    Args:
        slides: List of slide content
        diagrams: List of generated diagrams
        images: List of generated images
        template_path: Optional template file
        output_path: Output file path

    Returns:
        Path to generated presentation
    """
    # Load template or create new presentation
    if template_path and template_path.exists():
        prs = Presentation(str(template_path))
        logger.info(f"Loaded template: {template_path}")
    else:
        prs = Presentation()
        logger.info("Creating new presentation")

    # Set presentation to 16:9 widescreen format
    prs.slide_width = Inches(13.333)  # 16:9 ratio width
    prs.slide_height = Inches(7.5)  # 16:9 ratio height
    logger.info("Set presentation format to 16:9 widescreen")

    # Create lookup maps for diagrams and images
    diagram_map = {d.slide_number: d for d in diagrams}
    image_map = {i.slide_number: i for i in images}

    # Add slides
    for slide_content in slides:
        try:
            _add_slide(prs, slide_content, diagram_map, image_map)
            logger.debug(f"Added slide {slide_content.slide_number}: {slide_content.title}")
        except Exception as e:
            logger.error(f"Failed to add slide {slide_content.slide_number}: {e}")
            # Continue with other slides

    # Save presentation
    prs.save(str(output_path))
    logger.info(f"Saved presentation to {output_path}")

    return output_path


def _add_slide(prs, content: SlideContent, diagrams: dict[int, SlideDiagram], images: dict[int, SlideImage]):
    """Add slide using proper placeholder system."""

    # Determine if we have visuals for this slide
    has_diagram = False
    if content.slide_number in diagrams:
        diagram_path = diagrams[content.slide_number].image_path
        has_diagram = diagram_path is not None and diagram_path.exists()

    has_image = False
    if content.slide_number in images:
        image_path = images[content.slide_number].image_path
        has_image = image_path is not None and image_path.exists()

    has_visual = has_diagram or has_image

    # Select appropriate layout based on content and visuals
    if content.slide_type.value == "title":
        layout_idx = 0  # Title Slide
    elif has_visual and content.bullet_points:
        # Use "Two Content" for text + visual side-by-side
        layout_idx = 3
    elif has_visual and not content.bullet_points:
        # Use "Picture with Caption" for visual-only
        layout_idx = 8 if len(prs.slide_layouts) > 8 else 5
    else:
        # Standard "Title and Content"
        layout_idx = 1

    # Get the layout and add slide
    try:
        layout = prs.slide_layouts[layout_idx]
    except IndexError:
        # Fallback to Title and Content or blank
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        layout_idx = 1 if len(prs.slide_layouts) > 1 else 0

    slide = prs.slides.add_slide(layout)

    # Set title using placeholder index 0 (always title)
    try:
        if 0 in slide.placeholders:
            slide.placeholders[0].text = content.title
        elif slide.shapes.title:
            slide.shapes.title.text = content.title
        else:
            # Manual title if no placeholder
            _add_manual_title(slide, content.title)
    except (KeyError, AttributeError):
        _add_manual_title(slide, content.title)

    # Add subtitle for title slides
    if content.slide_type.value == "title" and content.subtitle:
        try:
            if 1 in slide.placeholders:
                slide.placeholders[1].text = content.subtitle
        except (KeyError, AttributeError):
            pass

    # Add bullet points if present
    if content.bullet_points:
        # For Two Content layout, text goes in placeholder 1
        # For Title and Content, text also goes in placeholder 1
        _fill_content_placeholder(slide, content.bullet_points)

    # Add visual if present
    if has_visual:
        diagram = diagrams.get(content.slide_number) if has_diagram else None
        image = images.get(content.slide_number) if has_image else None
        _add_visual_to_placeholder(slide, diagram, image, layout_idx, prs)

    # Add speaker notes
    if content.speaker_notes:
        slide.notes_slide.notes_text_frame.text = content.speaker_notes


def _add_manual_title(slide, title: str):
    """Add title manually when no placeholder exists."""
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_box.text_frame.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True


def _fill_content_placeholder(slide, bullet_points: list[str]):
    """Fill content placeholder with bullets."""
    if not bullet_points:
        return

    # Try to use placeholder 1 (content placeholder)
    try:
        if 1 in slide.placeholders:
            placeholder = slide.placeholders[1]
            text_frame = placeholder.text_frame
            text_frame.clear()

            for i, bullet in enumerate(bullet_points):
                clean_text = bullet.lstrip("•-*").strip()
                if i == 0:
                    text_frame.text = clean_text
                else:
                    p = text_frame.add_paragraph()
                    p.text = clean_text
                    p.level = 0
            return
    except (KeyError, AttributeError):
        pass

    # Fallback: find any content placeholder
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            text_frame = shape.text_frame
            text_frame.clear()

            for i, bullet in enumerate(bullet_points):
                clean_text = bullet.lstrip("•-*").strip()
                if i == 0:
                    text_frame.text = clean_text
                else:
                    p = text_frame.add_paragraph()
                    p.text = clean_text
                    p.level = 0
            return

    # Last resort: create new text box
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(6)  # Narrower to leave room for visuals
    height = Inches(5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame

    for i, bullet in enumerate(bullet_points):
        clean_text = bullet.lstrip("•-*").strip()
        if i == 0:
            text_frame.text = f"• {clean_text}"
        else:
            p = text_frame.add_paragraph()
            p.text = f"• {clean_text}"


def _add_visual_to_placeholder(slide, diagram: SlideDiagram | None, image: SlideImage | None, layout_idx: int, prs):
    """Add visual using placeholder if available, else manual position."""

    # Determine which visual to use
    visual_path = None
    if diagram and diagram.image_path and diagram.image_path.exists():
        visual_path = diagram.image_path
    elif image and image.image_path and image.image_path.exists():
        visual_path = image.image_path

    if not visual_path:
        return

    # Look for picture placeholder
    picture_ph = None
    for placeholder in slide.placeholders:
        try:
            if (
                hasattr(placeholder, "placeholder_format")
                and placeholder.placeholder_format.type == PP_PLACEHOLDER.PICTURE
            ):
                picture_ph = placeholder
                break
        except AttributeError:
            continue

    # Use placeholder if available
    if picture_ph:
        try:
            # PROPER WAY: Use insert_picture() - respects placeholder bounds!
            picture_ph.insert_picture(str(visual_path))
            logger.debug("Inserted visual using picture placeholder")
            return
        except Exception as e:
            logger.debug(f"Failed to use picture placeholder: {e}")

    # For Two Content layout (idx 3), try to use placeholder 2 (right content)
    if layout_idx == 3:
        try:
            # Find the right content placeholder (format.idx == 2)
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.idx == 2:
                    _manual_image_in_placeholder(slide, placeholder, visual_path)
                    logger.debug("Added visual to right content placeholder")
                    return
        except (KeyError, AttributeError) as e:
            logger.debug(f"Could not use placeholder 2: {e}")

    # Fallback: manual positioning
    _manual_positioning(slide, visual_path, prs)


def _manual_image_in_placeholder(slide, placeholder, image_path):
    """Add image within placeholder bounds."""
    # Get placeholder position and size
    left = placeholder.left
    top = placeholder.top
    width = placeholder.width
    height = placeholder.height

    # Add picture within these bounds
    pic = slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)

    # Maintain aspect ratio
    aspect_ratio = pic.width / pic.height
    if pic.width > width:
        pic.width = width
        pic.height = int(width / aspect_ratio)
    if pic.height > height:
        pic.height = height
        pic.width = int(height * aspect_ratio)


def _manual_positioning(slide, visual_path, prs):
    """Fallback manual positioning when no placeholder available."""
    # Check if slide has text content
    has_content = False
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title and shape.text_frame.text.strip():
            has_content = True
            break

    if has_content:
        # Position on right side
        left = Inches(7)
        top = Inches(1.5)
        max_width = Inches(6)
        max_height = Inches(5)

        pic = slide.shapes.add_picture(str(visual_path), left, top, height=max_height)

        # Adjust if too wide
        if pic.width > max_width:
            pic.width = max_width
    else:
        # Center the visual
        left = Inches(1.5)
        top = Inches(1.5)
        height = Inches(5.5)

        pic = slide.shapes.add_picture(str(visual_path), left, top, height=height)

        # Center horizontally
        slide_width = prs.slide_width
        pic.left = int((slide_width - pic.width) / 2)

    logger.debug("Added visual using manual positioning")


def _get_layout_index(content: SlideContent) -> int:
    """Determine the best layout index for the slide type.

    Note: For slides with visuals, we prefer layouts with content placeholders
    that can properly position text and images without overlap.
    """
    # Check if slide has visual content
    has_visual = content.visual_type.value != "none"

    layout_map = {
        "title": 0,  # Title slide
        "section": 2,  # Section header
        "bullets": 3 if has_visual else 1,  # Two Content layout if visual, else Title and Content
        "content": 3 if has_visual else 1,  # Two Content layout if visual, else Title and Content
        "diagram": 5,  # Blank for custom layout
        "image": 5,  # Blank for custom layout
        "mixed": 3,  # Two Content layout for mixed content
        "conclusion": 0,  # Title slide style
    }
    return layout_map.get(content.slide_type.value, 1)


def _add_title_slide_content(slide, content: SlideContent):
    """Add content specific to title slides."""
    # Title is already set
    # Add subtitle if exists
    if content.subtitle:
        try:
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = content.subtitle
        except (KeyError, IndexError):
            pass


def _add_bullet_slide_content(slide, content: SlideContent):
    """Add bullet points to slide."""
    if not content.bullet_points:
        return

    # Try to find content placeholder
    content_placeholder = None
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title:
            content_placeholder = shape
            break

    if content_placeholder:
        # Use existing placeholder
        text_frame = content_placeholder.text_frame
        text_frame.clear()  # Clear any existing text

        for i, bullet in enumerate(content.bullet_points):
            # Strip bullet characters from text (they'll be added by PowerPoint)
            clean_text = bullet.lstrip("•-*").strip()

            if i == 0:
                text_frame.text = clean_text
            else:
                p = text_frame.add_paragraph()
                p.text = clean_text
                p.level = 0
    else:
        # Create new text box for bullets
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(9)
        height = Inches(5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame

        for i, bullet in enumerate(content.bullet_points):
            # Strip bullet characters from text
            clean_text = bullet.lstrip("•-*").strip()

            if i == 0:
                text_frame.text = f"• {clean_text}"
            else:
                p = text_frame.add_paragraph()
                p.text = f"• {clean_text}"


def _add_standard_content(slide, content: SlideContent):
    """Add standard content to slide."""
    # Similar to bullet content but without bullet formatting
    if content.bullet_points:
        _add_bullet_slide_content(slide, content)


def _add_diagram(slide, diagram: SlideDiagram, prs):
    """Add diagram image to slide with smart positioning."""
    if not diagram.image_path or not diagram.image_path.exists():
        logger.warning(f"Diagram image not found: {diagram.image_path}")
        return

    try:
        # Check if slide has content (bullets/text) to determine positioning
        has_content = False
        content_width = Inches(6)  # Default content area width

        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title and shape.text_frame.text.strip():
                has_content = True
                # Get the width of the content area for better positioning
                if shape.width < prs.slide_width * 0.6:  # If content takes less than 60% width
                    content_width = shape.width
                break

        if has_content:
            # Position diagram on right side, avoiding overlap with text
            # Leave margin between text and image
            left = content_width + Inches(1)  # Position after content with margin
            top = Inches(1.5)
            max_width = prs.slide_width - left - Inches(0.5)  # Leave right margin
            max_height = Inches(5)

            # Calculate dimensions maintaining aspect ratio
            pic = slide.shapes.add_picture(str(diagram.image_path), left, top, height=max_height)

            # Adjust if image is too wide
            if pic.width > max_width:
                pic.width = max_width
                # Height will adjust automatically to maintain aspect ratio
        else:
            # Center diagram if no text content
            left = Inches(1.5)
            top = Inches(1.5)
            height = Inches(5.5)

            pic = slide.shapes.add_picture(str(diagram.image_path), left, top, height=height)

            # Center horizontally
            slide_width = prs.slide_width
            pic.left = int((slide_width - pic.width) / 2)

        logger.debug(f"Added diagram to slide {diagram.slide_number} ({'right side' if has_content else 'centered'})")
    except Exception as e:
        logger.error(f"Failed to add diagram: {e}")


def _add_image(slide, image: SlideImage, prs):
    """Add image to slide with smart positioning."""
    if not image.image_path or not image.image_path.exists():
        logger.warning(f"Image not found: {image.image_path}")
        return

    try:
        # Check if slide has content (bullets/text) to determine positioning
        has_content = False
        content_width = Inches(6)  # Default content area width

        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title and shape.text_frame.text.strip():
                has_content = True
                # Get the width of the content area for better positioning
                if shape.width < prs.slide_width * 0.6:  # If content takes less than 60% width
                    content_width = shape.width
                break

        if has_content:
            # Position image on right side, avoiding overlap with text
            # Leave margin between text and image
            left = content_width + Inches(1)  # Position after content with margin
            top = Inches(1.5)
            max_width = prs.slide_width - left - Inches(0.5)  # Leave right margin
            max_height = Inches(5)

            # Calculate dimensions maintaining aspect ratio
            pic = slide.shapes.add_picture(str(image.image_path), left, top, height=max_height)

            # Adjust if image is too wide
            if pic.width > max_width:
                pic.width = max_width
                # Height will adjust automatically to maintain aspect ratio
        else:
            # Center image if no text content
            left = Inches(1.5)
            top = Inches(1.5)
            height = Inches(5.5)

            pic = slide.shapes.add_picture(str(image.image_path), left, top, height=height)

            # Center horizontally
            slide_width = prs.slide_width
            pic.left = int((slide_width - pic.width) / 2)

        logger.debug(f"Added image to slide {image.slide_number} ({'right side' if has_content else 'centered'})")
    except Exception as e:
        logger.error(f"Failed to add image: {e}")
