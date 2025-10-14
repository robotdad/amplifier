"""Generate customizable PowerPoint template."""

import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger(__name__)


def generate_template(output_path: Path) -> Path:
    """Generate a default template with named placeholders.

    This template can be customized in PowerPoint and used with the generator.

    Returns:
        Path to generated template
    """
    prs = Presentation()

    # Set to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # The default template already has good layouts, but we'll document them
    # and add a sample slide showing the layout names

    # Add an instruction slide
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)

    slide.placeholders[0].text = "PPT Generator Template"  # type: ignore

    instructions = [
        "This template is ready to use with the PPT Generator",
        "Customize colors, fonts, and styling in PowerPoint",
        "Keep the same layout structure for compatibility",
        "",
        "Available Layouts:",
        "• Layout 0: Title Slide",
        "• Layout 1: Title and Content (text-only slides)",
        "• Layout 3: Two Content (text + visual side-by-side)",
        "• Layout 8: Picture with Caption (visual-focused slides)",
        "",
        "The generator automatically selects the best layout",
        "based on your content and visual needs.",
    ]

    text_frame = slide.placeholders[1].text_frame  # type: ignore
    text_frame.clear()
    for i, instruction in enumerate(instructions):
        if i == 0:
            text_frame.text = instruction
        else:
            p = text_frame.add_paragraph()
            p.text = instruction

    # Add example of each layout type

    # Example: Title Slide (Layout 0)
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    try:
        if 0 in slide.placeholders:  # type: ignore
            ph = slide.placeholders[0]
            ph.text = "Example: Title Slide"  # type: ignore
        if 1 in slide.placeholders:  # type: ignore
            ph = slide.placeholders[1]
            ph.text = "This is Layout 0 - Used for title and conclusion slides"  # type: ignore
    except (KeyError, AttributeError):
        pass

    # Example: Title and Content (Layout 1)
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    try:
        if 0 in slide.placeholders:  # type: ignore
            ph = slide.placeholders[0]
            ph.text = "Example: Title and Content"  # type: ignore
        if 1 in slide.placeholders:  # type: ignore
            ph = slide.placeholders[1]
            text_frame = ph.text_frame  # type: ignore
            text_frame.text = "This is Layout 1 - Used for text-only slides"
            p = text_frame.add_paragraph()
            p.text = "• Bullet point 1"
            p = text_frame.add_paragraph()
            p.text = "• Bullet point 2"
            p = text_frame.add_paragraph()
            p.text = "• Bullet point 3"
    except (KeyError, AttributeError):
        pass

    # Example: Two Content (Layout 3) if available
    if len(prs.slide_layouts) > 3:
        layout = prs.slide_layouts[3]
        slide = prs.slides.add_slide(layout)
        try:
            if 0 in slide.placeholders:  # type: ignore
                ph = slide.placeholders[0]
                ph.text = "Example: Two Content"  # type: ignore
            if 1 in slide.placeholders:  # type: ignore
                ph = slide.placeholders[1]
                text_frame = ph.text_frame  # type: ignore
                text_frame.text = "Left side: Text content"
                p = text_frame.add_paragraph()
                p.text = "• Used for slides with both text and visuals"
                p = text_frame.add_paragraph()
                p.text = "• Text goes on the left"
                p = text_frame.add_paragraph()
                p.text = "• Visual goes on the right"
            if 2 in slide.placeholders:  # type: ignore
                ph = slide.placeholders[2]
                ph.text = "Right side: Placeholder for images/diagrams"  # type: ignore
        except (KeyError, AttributeError):
            pass

    # Example: Picture with Caption (Layout 8) if available
    if len(prs.slide_layouts) > 8:
        layout = prs.slide_layouts[8]
        slide = prs.slides.add_slide(layout)
        try:
            if 0 in slide.placeholders:  # type: ignore
                ph = slide.placeholders[0]
                ph.text = "Example: Picture with Caption"  # type: ignore
            # Picture placeholder will be at index 1
            if 2 in slide.placeholders:  # type: ignore
                ph = slide.placeholders[2]
                ph.text = "This is Layout 8 - Used for visual-focused slides with captions"  # type: ignore
        except (KeyError, AttributeError):
            pass

    # Save template
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    logger.info(f"Generated template: {output_path}")
    logger.info("Customize this template in PowerPoint:")
    logger.info("  1. Open the template file")
    logger.info("  2. Go to View → Slide Master")
    logger.info("  3. Customize fonts, colors, logos")
    logger.info("  4. Save and use with --template flag")

    return output_path
