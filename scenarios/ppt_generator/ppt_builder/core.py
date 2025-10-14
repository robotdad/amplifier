"""Build PowerPoint presentation from generated content."""

import logging
from pathlib import Path

from PIL import Image
from pptx import Presentation

from ..models import SlideContent

logger = logging.getLogger(__name__)


def build_presentation(
    slides: list[SlideContent],
    output_path: Path,
    template_path: Path | None = None,
    images: dict[int, str] | None = None,
) -> Path:
    """Build PowerPoint presentation from slide content.

    Args:
        slides: List of slide content
        output_path: Output file path
        template_path: Optional template PPTX file
        images: Optional dict mapping slide numbers to image paths

    Returns:
        Path to created presentation
    """
    # Create presentation from template or blank
    if template_path and template_path.exists():
        prs = Presentation(str(template_path))
        logger.info(f"Using template: {template_path}")
    else:
        prs = Presentation()
        logger.info("Creating presentation from blank")

    # Clear any existing slides if using template
    # Note: We'll just add new slides, not remove template slides
    # as removing slides from templates can be complex

    # Add slides
    for slide_content in slides:
        _add_slide(prs, slide_content, images)

    # Save presentation
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    logger.info(f"Presentation saved to: {output_path}")

    return output_path


def _add_slide(prs, content: SlideContent, images: dict[int, str] | None = None):
    """Add a slide to the presentation.

    Args:
        prs: Presentation object
        content: Slide content
        images: Optional image paths
    """
    # Select appropriate layout
    layout = _get_layout(prs, content.layout_preference)
    slide = prs.slides.add_slide(layout)

    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = content.title

    # Handle different layouts
    if content.layout_preference == "title" or content.slide_type.value == "title":
        # Title slide with subtitle
        if content.subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = content.subtitle

    else:
        # Content slides
        # Find content placeholder
        content_placeholder = None
        for shape in slide.placeholders:
            if shape.has_text_frame and shape != slide.shapes.title:
                content_placeholder = shape
                break

        if content_placeholder and content.bullet_points:
            text_frame = content_placeholder.text_frame
            text_frame.clear()  # Clear any template text

            for i, point in enumerate(content.bullet_points):
                if i == 0:
                    # First point
                    p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
                    p.text = point
                else:
                    # Additional points
                    p = text_frame.add_paragraph()
                    p.text = point
                    p.level = 0  # Top-level bullet

        # Add subtitle if present and there's a suitable placeholder
        if content.subtitle:
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # Subtitle placeholder
                    shape.text = content.subtitle
                    break

    # Add image if available
    if images and content.slide_number in images:
        image_path = Path(images[content.slide_number])
        if image_path.exists():
            _add_image_to_slide(slide, image_path)

    # Add speaker notes
    if content.speaker_notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = content.speaker_notes


def _get_layout(prs, layout_name: str):
    """Get slide layout by name or index.

    Args:
        prs: Presentation object
        layout_name: Layout name

    Returns:
        Slide layout object
    """
    # Map our layout names to typical PPT layout indices
    layout_map = {
        "title_slide": 0,  # Title Slide
        "title_and_content": 1,  # Title and Content
        "two_content": 3,  # Two Content
        "comparison": 4,  # Comparison
        "blank": 6,  # Blank
    }

    # Try to get by index
    layout_idx = layout_map.get(layout_name, 1)  # Default to Title and Content

    # Ensure we don't exceed available layouts
    if layout_idx >= len(prs.slide_layouts):
        layout_idx = 1  # Fall back to Title and Content

    return prs.slide_layouts[layout_idx]


def _add_image_to_slide(slide, image_path: Path):
    """Add an image to a slide as background or in placeholder.

    Args:
        slide: Slide object
        image_path: Path to image file
    """
    try:
        # Try to add image in a picture placeholder if available
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 18:  # Picture placeholder
                shape.insert_picture(str(image_path))
                return

        # Otherwise, add as a positioned image
        # Get slide dimensions
        prs = slide.part.presentation
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Open image to get dimensions
        img = Image.open(image_path)
        img_width, img_height = img.size

        # Calculate position to center on right side of slide
        # Leave left 60% for content
        left = int(slide_width * 0.6)
        available_width = int(slide_width * 0.35)

        # Scale image to fit
        scale = min(available_width / img_width, slide_height / img_height) * 0.8
        width = int(img_width * scale)
        height = int(img_height * scale)

        # Center vertically
        top = int((slide_height - height) / 2)

        # Add the image
        slide.shapes.add_picture(str(image_path), left, top, width, height)

        logger.debug(f"Added image to slide: {image_path.name}")

    except Exception as e:
        logger.warning(f"Failed to add image to slide: {e}")
